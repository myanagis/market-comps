# pages/2_PDF_Parser.py
"""
PDF Parser — Upload and parse PDF documents using OpenRouter PDF plugins.
Currently supports: Term Sheet extraction and general document summarization.
"""
from __future__ import annotations

import streamlit as st

from market_comps.config import settings, MODEL_OPTIONS
from market_comps.pdf_parser import TermExtractor
from market_comps.pdf_parser.models import ParserResult
from market_comps.pdf_parser.pdf_client import PDF_ENGINE_PRICING
from streamlit_paste_button import paste_image_button
import io

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PDF Parser",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ── CSS ───────────────────────────────────────────────────────────────────────
from market_comps.ui import inject_global_style

inject_global_style()

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
.doc-type-badge {
    display: inline-block;
    padding: 0.35rem 0.8rem;
    border-radius: 999px;
    font-size: 0.82rem;
    font-weight: 600;
    margin-bottom: 0.3rem;
}
.badge-ts   { background: #14532d; color: #86efac; border: 1px solid #166534; }
.badge-safe { background: #1e3a5f; color: #93c5fd; border: 1px solid #1d4ed8; }
.badge-conv { background: #3b1f0e; color: #fcd34d; border: 1px solid #92400e; }
.badge-loi  { background: #2e1065; color: #c4b5fd; border: 1px solid #6d28d9; }
.badge-unk  { background: #1e293b; color: #94a3b8; border: 1px solid #475569; }

.rationale-text { color: #64748b; font-size: 0.83rem; font-style: italic; margin-top: 0.2rem; }

.usage-badge {
    background: #0f2744; border: 1px solid #1e4a7a; border-radius: 8px;
    padding: 0.6rem 0.8rem; font-size: 0.78rem; color: #93c5fd; margin-top: 0.8rem;
}
.usage-badge b { color: #bfdbfe; }

.info-box {
    background: #1e293b; border-left: 4px solid #34d399;
    border-radius: 0 8px 8px 0; padding: 0.8rem 1rem;
    color: #94a3b8; font-size: 0.9rem;
}

/* Quote chips */
.quote-chip {
    display: block;
    background: #0f172a; border: 1px solid #1e3a5f;
    border-radius: 6px; padding: 0.3rem 0.6rem;
    font-size: 0.78rem; color: #7dd3fc;
    font-style: italic; margin-top: 0.25rem;
    white-space: pre-wrap; word-break: break-word;
}
.snippet-chip {
    display: block;
    background: #1c1917; border: 1px solid #44403c;
    border-radius: 6px; padding: 0.3rem 0.6rem;
    font-size: 0.78rem; color: #a8a29e;
    font-style: italic; margin-top: 0.25rem;
    white-space: pre-wrap; word-break: break-word;
}
</style>
""", unsafe_allow_html=True)

# ── Header ────────────────────────────────────────────────────────────────────
st.markdown("""
<h1>📄 Document Parser</h1>
<p>Upload a PDF or PPTX to extract key terms (term sheets), generate a summary, or transcribe visually.</p>
""", unsafe_allow_html=True)

# ── Session state ─────────────────────────────────────────────────────────────
if "pdf_result" not in st.session_state:
    st.session_state["pdf_result"] = None

# ── Engine options ────────────────────────────────────────────────────────────
ENGINE_OPTIONS = {
    "PDF Text (Free)": "pdf-text",
    "Mistral OCR ($2 / 1k pages)": "mistral-ocr",
    "Native (input tokens)": "native",
}

# ── Upload ────────────────────────────────────────────────────────────────────
col1, col2 = st.columns([2, 1])
with col1:
    uploaded_file = st.file_uploader(
        "Upload a PDF, PPTX, or Image",
        type=["pdf", "pptx", "png", "jpg", "jpeg"],
        label_visibility="collapsed",
    )

with col2:
    st.markdown("<div style='margin-top: 5px; margin-bottom: 5px; color: #888;'>Or paste an image:</div>", unsafe_allow_html=True)
    paste_result = paste_image_button(
        label="📋 Paste Image",
        background_color="#FF4B4B",
        hover_background_color="#FF0000",
        errors='ignore'
    )

file_bytes = None
filename = ""
file_size_kb = 0

if uploaded_file is not None:
    file_bytes = uploaded_file.read()
    filename = uploaded_file.name
    file_size_kb = uploaded_file.size / 1024
elif paste_result.image_data is not None:
    img_byte_arr = io.BytesIO()
    paste_result.image_data.save(img_byte_arr, format='PNG')
    file_bytes = img_byte_arr.getvalue()
    filename = "pasted_image.png"
    file_size_kb = len(file_bytes) / 1024

if file_bytes is None:
    st.markdown("""
    <div class="info-box">
        👆 Upload a <b>PDF, PPTX, or Image</b> or paste an image above to get started — term sheets, SAFE notes,
        convertible notes, or any other document.
    </div>
    """, unsafe_allow_html=True)
    st.session_state["pdf_result"] = None
else:
    st.caption(f"📎 **{filename}** — {file_size_kb:.1f} KB")

    action = st.radio(
        "Action",
        ["Extract Term Sheet (Structured)", "Transcribe Pitch Deck (High-Fidelity)"],
        horizontal=True
    )

    # Advanced Options
    def format_model(m: str) -> str:
        in_price, out_price = settings.get_model_pricing(m)
        return f"{m} (${in_price:.2f} / ${out_price:.2f})"

    with st.expander("⚙️ Advanced Options", expanded=False):
        _ao1, _ao2 = st.columns([3, 1])
        with _ao1:
            model = st.selectbox(
                "LLM Model (Prices shown: $input / $output per 1M tokens)",
                MODEL_OPTIONS,
                index=MODEL_OPTIONS.index(settings.default_model)
                if settings.default_model in MODEL_OPTIONS else 0,
                format_func=format_model,
            )
        with _ao2:
            if action == "Extract Term Sheet (Structured)":
                ENGINE_OPTIONS = {
                    "PDF Text (Free)": "pdf-text",
                    "Mistral OCR ($2 / 1k pages)": "mistral-ocr",
                    "Native (input tokens)": "native",
                }
                engine_label = st.selectbox("PDF Engine", list(ENGINE_OPTIONS.keys()), index=0)
                engine = ENGINE_OPTIONS[engine_label]
            else:
                METHOD_OPTIONS = {
                    "OCR (Mistral)": "ocr",
                    "Text Reader (Native)": "text",
                    "VLM (Image-based)": "vlm",
                    "Hybrid (VLM + Text cross-check)": "vlm_plus_text"
                }
                method_label = st.selectbox("Extraction Method", list(METHOD_OPTIONS.keys()), index=3)
                method = METHOD_OPTIONS[method_label]

    parse_clicked = st.button("🔍 Parse Document", type="primary")

    # ── Run pipeline ──────────────────────────────────────────────────────────
    if parse_clicked:
        file_bytes = uploaded_file.read()
        filename = uploaded_file.name

        # --- PPTX Flow ---
        if filename.lower().endswith(".pptx"):
            with st.spinner("🔄 Extracting text from presentation..."):
                import io
                from pptx import Presentation
                
                prs = Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        text_runs.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
                
                extracted_text = "\n\n".join(text_runs)
                
                class PPTXResult:
                    def __init__(self):
                        self.document_type = "presentation"
                        self.doc_type_confidence = "high"
                        self.doc_type_rationale = "Directly extracted from PPTX file."
                        from market_comps.models import LLMUsage
                        self.llm_usage = LLMUsage()
                        self.model_used = "N/A"
                        self.pdf_engine = "python-pptx"
                        self.summary = extracted_text
                        self.raw_extracted_text = extracted_text
                        self.terms = []
                        self.errors = []
                        
                st.session_state["pdf_result"] = PPTXResult()

        # --- PDF Flow ---
        else:
            progress = st.empty()
    
            with progress.container():
                st.info("🔄 **Step 0** — Processing document…")
    
            try:
                if action == "Transcribe Pitch Deck (High-Fidelity)":
                    from market_comps.pdf_parser.document_processor import process_document
                    
                    with st.spinner(f"Transcribing via {method}..."):
                        content, usage = process_document(file_bytes, filename, method, model)
                        
                        class TranscriptionResult:
                            def __init__(self):
                                self.document_type = "presentation"
                                self.doc_type_confidence = "high"
                                self.doc_type_rationale = f"Transcribed via {method}"
                                self.llm_usage = usage
                                self.model_used = model
                                self.pdf_engine = method
                                self.summary = content
                                self.raw_extracted_text = content
                                self.terms = []
                                self.errors = []
                        
                        st.session_state["pdf_result"] = TranscriptionResult()
                        progress.empty()

                else:
                    extractor = TermExtractor(model=model, pdf_engine=engine)
                    with st.spinner(""):
                        result = extractor.run(pdf_bytes=file_bytes, filename=filename)
        
                    st.session_state["pdf_result"] = result
                    progress.empty()
    
            except Exception as exc:
                progress.empty()
                st.error(f"❌ Parser error: {exc}")
                st.session_state["pdf_result"] = None

# ── Results ───────────────────────────────────────────────────────────────────
result = st.session_state.get("pdf_result")

# Normalise a quote item that may be a plain str (old session state)
# or a SupportingQuote dataclass (new format).
def _as_quote(q):
    if isinstance(q, str):
        class _Q:
            text = q
            page = None
        return _Q()
    return q

if result is not None:

    # Errors
    for err in result.errors:
        st.warning(f"⚠️ {err}")

    # ── Document Type badge ────────────────────────────────────────────────
    BADGE_CLASS = {
        "term_sheet": "badge-ts",
        "safe_note": "badge-safe",
        "convertible_note": "badge-conv",
        "loi": "badge-loi",
        "letter_of_intent": "badge-loi",
        "presentation": "badge-loi",
        "pdf": "badge-loi",
    }
    LABEL_MAP = {
        "term_sheet": "📋 Term Sheet",
        "safe_note": "📄 SAFE Note",
        "convertible_note": "📝 Convertible Note",
        "loi": "🤝 Letter of Intent",
        "letter_of_intent": "🤝 Letter of Intent",
        "presentation": "📊 PowerPoint Presentation",
        "pdf": "📄 PDF Document (VLM)",
        "other": "📂 Other Document",
    }
    badge_cls = BADGE_CLASS.get(result.document_type, "badge-unk")
    label = LABEL_MAP.get(result.document_type, f"📂 {result.document_type.replace('_', ' ').title()}")

    st.markdown(f"""
    <div class="section-header">Document Classification</div>
    <span class="doc-type-badge {badge_cls}">{label}</span>
    <span style="color:#64748b; font-size:0.82rem; margin-left:0.5rem;">
        Confidence: <b>{result.doc_type_confidence}</b>
    </span>
    <div class="rationale-text">{result.doc_type_rationale}</div>
    """, unsafe_allow_html=True)

    # ── Usage & cost breakdown ─────────────────────────────────────────────
    u = result.llm_usage
    pages = getattr(result, "pdf_pages", 0)
    engine = getattr(result, "pdf_engine", result.pdf_engine)
    is_ocr = engine == "mistral-ocr"

    _m1, _m2, _m3 = st.columns(3)
    _m1.metric("API Calls", u.call_count)
    _m2.metric("Tokens", f"{u.total_tokens:,}")
    _m3.metric(
        "Total Cost",
        f"${u.estimated_cost_usd:.5f}",
        help=("OpenRouter usage cost — includes Mistral-OCR PDF parsing fee."
              if is_ocr else "OpenRouter LLM usage cost. pdf-text engine is free."),
    )
    caption = f"Model: `{result.model_used}` · Engine: `{engine}`"
    if pages:
        caption += f" · {pages} pages"
    if is_ocr:
        caption += " · PDF parsing cost included in Total Cost"
    st.caption(caption)


    # ── Term extraction results ────────────────────────────────────────────
    if result.terms:
        st.markdown('<div class="section-header">📊 Extracted Terms</div>', unsafe_allow_html=True)

        # Summary table first
        import pandas as pd
        rows = []
        for term in result.terms:
            conf_label = "N/A" if term.confidence == "not_found" else term.confidence
            first_quote = ""
            if term.supporting_quotes:
                q = _as_quote(term.supporting_quotes[0])
                first_quote = f'"{q.text[:120]}"'
                if q.page:
                    first_quote += f" (p.{q.page})"
            elif term.possible_snippets:
                first_quote = "~" + term.possible_snippets[0][:100]
            rows.append({
                "Field": term.name,
                "Value": term.value or "—",
                "Confidence": conf_label,
                "Supporting Quote": first_quote,
            })
        st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

        # Detailed expandable rows below
        st.markdown('<div class="section-header">📋 Details</div>', unsafe_allow_html=True)
        for term in result.terms:
            conf_label = "N/A" if term.confidence == "not_found" else term.confidence
            conf_icon = {"high": "✅", "low": "⚠️", "not_found": "—"}.get(term.confidence, "—")
            value_display = term.value if term.value else "—"

            with st.expander(f"{conf_icon} **{term.name}** — {value_display}", expanded=(term.confidence == "high")):
                cols = st.columns([1, 2])
                with cols[0]:
                    st.markdown(f"**Confidence:** {conf_label}")
                    st.markdown(f"**Value:** {value_display}")

                with cols[1]:
                    if term.supporting_quotes:
                        st.markdown("**Supporting quotes:**")
                        for q_raw in term.supporting_quotes:
                            q = _as_quote(q_raw)
                            page_badge = f' <span style="color:#64748b; font-size:0.72rem;">p.{q.page}</span>' if q.page else ""
                            preview = q.text[:160] + ("…" if len(q.text) > 160 else "")
                            escaped = q.text.replace('"', '&quot;').replace("'", "&#39;")
                            st.markdown(
                                f'<span class="quote-chip" title="{escaped}">"{preview}"{page_badge}</span>',
                                unsafe_allow_html=True,
                            )
                    elif term.possible_snippets:
                        st.markdown("**Nearby snippets (uncertain):**")
                        for s in term.possible_snippets:
                            escaped = s.replace('"', '&quot;').replace("'", "&#39;")
                            preview = s[:160] + ("…" if len(s) > 160 else "")
                            st.markdown(
                                f'<span class="snippet-chip" title="{escaped}">{preview}</span>',
                                unsafe_allow_html=True,
                            )
                    else:
                        st.markdown('<span style="color:#475569; font-size:0.85rem;">No relevant text found.</span>', unsafe_allow_html=True)

    # ── Summary results ────────────────────────────────────────────────────
    elif result.summary:
        st.markdown('<div class="section-header">📝 Document Summary</div>', unsafe_allow_html=True)
        st.markdown(result.summary)

    # ── Troubleshooting: raw extracted text ───────────────────────────────
    with st.expander("🔍 Troubleshooting — Raw Extracted Text", expanded=False):
        # Debug metadata
        import dataclasses, json as _json
        debug_info = {
            "pdf_engine": getattr(result, "pdf_engine", "?"),
            "pdf_pages": getattr(result, "pdf_pages", "NOT SET"),
            "pdf_parsing_cost_usd": getattr(result, "pdf_parsing_cost_usd", "NOT SET"),
            "llm_estimated_cost_usd": result.llm_usage.estimated_cost_usd,
            "document_type": result.document_type,
        }
        st.json(debug_info)
        st.divider()
        raw_text = getattr(result, "raw_extracted_text", None)
        if raw_text:
            st.code(raw_text, language=None)
        else:
            st.caption("No raw text captured. Re-run the parser to populate this field.")


# ── How It Works ──────────────────────────────────────────────────────────────
with st.expander("🤖 How this page works", expanded=False):
    st.markdown("""
**PDF Parser** extracts key legal and financial terms from uploaded documents using a two-step LLM pipeline:

1. **PDF Parsing** — The document is sent to OpenRouter's file-parser plugin (supports `mistral-ocr‑2503`
   and `gemini-2.5-flash` engines) which converts the PDF into structured text with page annotations.

2. **Term Extraction** — A second LLM call reads the extracted text and populates up to 21 structured fields
   (dates, parties, valuations, rights, provisions, etc.) with verbatim quotes and page-number references.

3. **Summarization** — A third call creates a plain-English executive summary.

Anti-hallucination rules are applied at each step: the model must quote verbatim from the document and
default to "not found" when information is absent. A troubleshooting section shows the raw extracted text
for verification.
""")
    
    st.markdown("#### LLM Prompts Used")
    from market_comps.pdf_parser.term_extractor import _CLASSIFY_PROMPT, _EXTRACT_PROMPT, _SUMMARIZE_PROMPT
    
    st.markdown("**1. Document Classification**")
    st.code(_CLASSIFY_PROMPT, language="text")
    
    st.markdown("**2. Term Extraction (for Term Sheets & SAFEs)**")
    st.code(_EXTRACT_PROMPT, language="text")
    
    st.markdown("**3. Summarization (for other document types)**")
    st.code(_SUMMARIZE_PROMPT, language="text")






