import time
import streamlit as st

st.set_page_config(
    page_title="Schema-Driven Framework",
    page_icon="📋",
    layout="wide",
    initial_sidebar_state="expanded",
)

from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from market_comps.pdf_parser.pdf_client import PDFClient

from market_comps.ui import inject_global_style
from market_comps.config import MODEL_OPTIONS
from market_comps.schema_framework.engine import run_schema_extraction, synthesize_evidence
from market_comps.config import settings

inject_global_style()

_DIR = Path(__file__).resolve().parent.parent / "market_comps" / "schema_framework"
_SCHEMA_PATH = _DIR / "starter_schema.json"
_PROMPT_PATH = _DIR / "evidence_extraction_prompt.md"
_SYNTH_PROMPT_PATH = _DIR / "synthesis_prompt.md"

def _load_asset(path: Path) -> str:
    with path.open("r", encoding="utf-8") as f:
        return f.read()

st.markdown("""
<h1>📋 Schema-Driven Framework</h1>
<p>Ingest and extract structured evidence based on a standardized framework, merging data and synthesizing insights.</p>
""", unsafe_allow_html=True)

# ── Session state ─────────────────────────────────────────────────────────────
if "sdf_sources" not in st.session_state:
    st.session_state["sdf_sources"] = [{"source_name": "article_1", "source_date": "", "text": ""}]
if "sdf_results" not in st.session_state:
    st.session_state["sdf_results"] = None
if "sdf_schema" not in st.session_state:
    st.session_state["sdf_schema"] = _load_asset(_SCHEMA_PATH)
if "sdf_ext_prompt" not in st.session_state:
    st.session_state["sdf_ext_prompt"] = _load_asset(_PROMPT_PATH)
if "sdf_synth_prompt" not in st.session_state:
    st.session_state["sdf_synth_prompt"] = _load_asset(_SYNTH_PROMPT_PATH)

# Configure model
def format_model(m: str) -> str:
    in_price, out_price = settings.get_model_pricing(m)
    return f"{m} (${in_price:.2f} / ${out_price:.2f})"

with st.expander("⚙️ Configuration", expanded=True):
    col1, _ = st.columns([1, 1])
    with col1:
        # Defaulting to a cheaper model (e.g. gpt-4o-mini or gemini-flash)
        default_idx = 0
        if "google/gemini-2.5-flash-lite" in MODEL_OPTIONS:
            default_idx = MODEL_OPTIONS.index("google/gemini-2.5-flash-lite")
        elif "openai/gpt-4o-mini" in MODEL_OPTIONS:
            default_idx = MODEL_OPTIONS.index("openai/gpt-4o-mini")
        
        selected_model = st.selectbox(
            "Extraction & Synthesis Model",
            options=MODEL_OPTIONS,
            index=default_idx,
            format_func=format_model,
        )

with st.expander("🛠️ Advanced Configuration", expanded=False):
    st.markdown("**Editable Schema & Prompts**")
    st.session_state["sdf_schema"] = st.text_area(
        "JSON Schema", 
        value=st.session_state["sdf_schema"], 
        height=250,
        help="Edit the JSON schema used for evidence extraction."
    )
    
    col_ext, col_syn = st.columns(2)
    with col_ext:
        st.session_state["sdf_ext_prompt"] = st.text_area(
            "Extraction Prompt", 
            value=st.session_state["sdf_ext_prompt"], 
            height=250,
            help="Prompt used for structured extraction from raw text."
        )
    with col_syn:
        st.session_state["sdf_synth_prompt"] = st.text_area(
            "Synthesis Prompt", 
            value=st.session_state["sdf_synth_prompt"], 
            height=250,
            help="Prompt used to synthesize all extracted evidence."
        )

# ── Inputs ────────────────────────────────────────────────────────────────────
st.markdown("### Raw Documents")

sources = st.session_state["sdf_sources"]

to_remove = None
for i, entry in enumerate(sources):
    st.markdown(f"**Source {i+1}**")
    c1, c1b, c2, c3 = st.columns([2, 2, 7, 1])
    with c1:
        sources[i]["source_name"] = st.text_input(
            "Source Identifier", value=entry.get("source_name", ""), key=f"src_name_{i}"
        )
    with c1b:
        sources[i]["source_date"] = st.text_input(
            "Date (Optional)", value=entry.get("source_date", ""), key=f"src_date_{i}"
        )
    with c2:
        sources[i]["text"] = st.text_area(
            "Raw Text", value=entry.get("text", ""), key=f"src_text_{i}", height=100
        )
    with c3:
        st.markdown("<div style='height: 28px'></div>", unsafe_allow_html=True) # align button
        if len(sources) > 1:
            if st.button("✕", key=f"src_rm_{i}", help="Remove"):
                to_remove = i

if to_remove is not None:
    sources.pop(to_remove)
    st.session_state["sdf_sources"] = sources
    st.rerun()

col_add, _ = st.columns([2, 10])
with col_add:
    if st.button("➕ Add raw document", type="secondary"):
        sources.append({"source_name": f"article_{len(sources)+1}", "source_date": "", "text": ""})
        st.session_state["sdf_sources"] = sources
        st.rerun()

st.markdown("---")
st.markdown("#### Or Upload a Document")

uploaded_file = st.file_uploader("Upload File", type=["txt", "docx", "pptx", "pdf"])

if uploaded_file is not None:
    filename = uploaded_file.name
    ext = filename.split(".")[-1].lower()
    file_bytes = uploaded_file.getvalue()
    
    if st.session_state.get("last_uploaded_file") != uploaded_file.file_id:
        extracted_text = ""
        
        if ext == "txt":
            extracted_text = file_bytes.decode("utf-8", errors="replace")
        elif ext == "docx":
            import io
            doc = Document(io.BytesIO(file_bytes))
            extracted_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif ext == "pptx":
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            extracted_text = "\n".join(text_runs)
        elif ext == "pdf":
            import io
            pdf = PdfReader(io.BytesIO(file_bytes))
            num_pages = len(pdf.pages)
            
            confirm_key = f"confirm_pdf_{uploaded_file.file_id}"
            
            if num_pages > 20 and not st.session_state.get(confirm_key, False):
                st.warning(f"This PDF has {num_pages} pages. Mistral OCR costs roughly $2 / 1000 pages. Do you want to proceed with extraction?")
                if st.button("Confirm Mistral OCR Processing"):
                    st.session_state[confirm_key] = True
                    st.rerun()
                else:
                    st.stop()
            
            with st.spinner("Extracting text via Mistral OCR... This might take a minute."):
                client = PDFClient(pdf_engine="mistral-ocr", model=selected_model)
                raw_text, _, usage = client.send(
                    prompt="Extract the raw text from this document as accurately as possible.", 
                    pdf_bytes=file_bytes, 
                    filename=filename
                )
                extracted_text = raw_text
                st.success(f"Extracted {num_pages} pages via Mistral OCR.")
        
        if extracted_text and ext != "pdf":
            with st.spinner("Processing file..."):
                pass # Already processed synchronously
                
        if extracted_text:
            sources.append({"source_name": filename, "source_date": "", "text": extracted_text})
            st.session_state["sdf_sources"] = sources
            st.session_state["last_uploaded_file"] = uploaded_file.file_id
            st.rerun()

valid_sources = [s for s in sources if s["text"].strip()]

# ── Run ───────────────────────────────────────────────────────────────────────
st.markdown("---")
_rb1, _ = st.columns([3, 10])
with _rb1:
    analyze_clicked = st.button(
        "🔍 Analyze & Extract",
        type="primary",
        disabled=not valid_sources,
        use_container_width=True,
    )

if analyze_clicked and valid_sources:
    st.session_state["sdf_results"] = None
    
    extracted_data = []
    
    with st.status("Processing Documents...", expanded=True) as status:
        for s in valid_sources:
            src_name = s.get("source_name", "Unknown")
            st.write(f"Extracting from **{src_name}**...")
            t0 = time.time()
            result = run_schema_extraction(
                text=s["text"],
                source_name=src_name,
                source_date=s.get("source_date", ""),
                model=selected_model,
                schema_text=st.session_state["sdf_schema"],
                extraction_prompt_template=st.session_state["sdf_ext_prompt"]
            )
            elapsed = time.time() - t0
            result["extraction_time_s"] = elapsed
            st.write(f"✅ Extracted from **{src_name}** in {elapsed:.1f}s")
            extracted_data.append(result)
        
        st.write("Synthesizing extracted evidence...")
        t_synth = time.time()
        synthesis, synth_usage = synthesize_evidence(
            extracted_data, 
            model=selected_model,
            synth_prompt_template=st.session_state["sdf_synth_prompt"]
        )
        synth_elapsed = time.time() - t_synth
        st.write(f"✅ Synthesized evidence in {synth_elapsed:.1f}s")
        
        status.update(label="Analysis Complete", state="complete")
        st.session_state["sdf_results"] = {
            "synthesis": synthesis,
            "evidence": extracted_data,
            "synth_usage": synth_usage,
            "synth_time_s": synth_elapsed
        }

# ── Results ───────────────────────────────────────────────────────────────────
results = st.session_state.get("sdf_results")
if results:
    st.markdown("## Results")

    tab1, tab2 = st.tabs(["Integrated Synthesis", "Structured Evidence Data"])
    
    with tab1:
        st.markdown("### Synthesis")
        st.info(results["synthesis"])

        su = results.get("synth_usage")
        synth_time = results.get("synth_time_s", 0.0)
        if su:
            st.markdown(
                f'<div style="font-size: 0.8rem; color: #64748b; margin-top: 10px;">'
                f'<b>Synthesis Cost:</b> {su.total_tokens:,} tokens | Est. ${su.estimated_cost_usd:.5f} | Time: {synth_time:.1f}s'
                f'</div>', unsafe_allow_html=True
            )

    with tab2:
        st.markdown("### Raw Extracted Records")
        for res in results["evidence"]:
            u = res.get("usage")
            ext_time = res.get("extraction_time_s", 0.0)
            usage_str = f" | {u.total_tokens:,} tokens | Est. ${u.estimated_cost_usd:.5f} | Time: {ext_time:.1f}s" if u else " | Usage unknown"
            with st.expander(f"📦 Evidence from {res['source']}{usage_str}", expanded=True):
                # Print Metadata
                meta_cols = st.columns(3)
                with meta_cols[0]:
                    st.markdown(f"**Source:** {res.get('source', 'Unknown')}")
                with meta_cols[1]:
                    if res.get('date'):
                        st.markdown(f"**Date:** {res.get('date')}")
                
                st.markdown("---")
                
                # Print Data structure
                data = res.get("data", {})
                if isinstance(data, dict) and "error" in data:
                    st.error(data["error"])
                elif isinstance(data, dict):
                    for category, subcategories in data.items():
                        if not isinstance(subcategories, dict):
                            continue
                            
                        # Only show category if it has at least one populated subcategory
                        has_items = any(bool(items) for items in subcategories.values())
                        if not has_items:
                            continue
                            
                        st.markdown(f"#### {category}")
                        
                        for sub_label, items in subcategories.items():
                            if items and isinstance(items, list):
                                st.markdown(f"**{sub_label}**")
                                for item in items:
                                    if isinstance(item, dict):
                                        # Try to find the actual quote text vs metadata tags
                                        text_val = item.get("quote", item.get("text", item.get("content", "")))
                                        if not text_val:
                                            st.markdown(f"- {item}")
                                        else:
                                            # Format extra metadata cleanly (tags, confidence, etc)
                                            extras = [f"{k}: {v}" for k, v in item.items() if k not in ["quote", "text", "content"]]
                                            extras_str = f" _({', '.join(extras)})_" if extras else ""
                                            st.markdown(f"- \"{text_val}\"{extras_str}")
                                    else:
                                        st.markdown(f"- {item}")
                        st.markdown("<br>", unsafe_allow_html=True)
                else:
                    st.json(data)
