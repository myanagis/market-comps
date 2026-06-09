import base64
import logging
from pathlib import Path
from typing import Tuple

import fitz  # PyMuPDF
from prefect import task, flow
from market_comps.llm_client import LLMClient
from market_comps.pdf_parser.pdf_client import PDFClient
from market_comps.models import LLMUsage

logger = logging.getLogger(__name__)

_DIR = Path(__file__).resolve().parent.parent / "pdf_parser"
_INSTRUCTIONS_PATH = _DIR / "vlm_instructions.md"

def _load_instructions() -> str:
    with _INSTRUCTIONS_PATH.open("r", encoding="utf-8") as f:
        return f.read()

@task(name="extract_text_with_vlm")
def extract_text_with_vlm(pdf_bytes: bytes, model: str, instructions: str) -> Tuple[str, LLMUsage]:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    messages = [{"role": "user", "content": [{"type": "text", "text": instructions}]}]
    
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("jpeg")
        b64 = base64.b64encode(img_bytes).decode("utf-8")
        messages[0]["content"].append({
            "type": "image_url", 
            "image_url": {"url": f"data:image/jpeg;base64,{b64}"}
        })
        
    client = LLMClient(model=model)
    return client.chat_completion(messages=messages, temperature=0.0, step_name="vlm_transcription")

@task(name="extract_text_with_ocr")
def extract_text_with_ocr(pdf_bytes: bytes, filename: str, model: str, instructions: str) -> Tuple[str, LLMUsage]:
    client = PDFClient(pdf_engine="mistral-ocr", model=model)
    content, _, step_usage = client.send(prompt=instructions, pdf_bytes=pdf_bytes, filename=filename, temperature=0.0, step_name="mistral_ocr")
    return content, step_usage

@task(name="extract_text_with_paddle_ocr")
def extract_text_with_paddle_ocr(pdf_bytes: bytes, filename: str) -> Tuple[str, LLMUsage]:
    try:
        from paddleocr import PaddleOCR
        import fitz
        import io
        import numpy as np
        from PIL import Image
        
        ocr = PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
        doc = fitz.open(stream=pdf_bytes, filetype="pdf" if filename.lower().endswith(".pdf") else None)
        
        full_text = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img_np = np.array(img)
            
            result = ocr.ocr(img_np, cls=True)
            page_text = []
            if result and result[0]:
                for line in result[0]:
                    page_text.append(line[1][0])
            full_text.append(f"===Page {page_num + 1}===\n" + "\n".join(page_text))
            
        return "\n\n".join(full_text), LLMUsage()
    except Exception as e:
        logger.error(f"PaddleOCR error: {e}")
        return f"PaddleOCR extraction failed: {e}", LLMUsage()

@task(name="extract_text_with_native")
def extract_text_with_native(pdf_bytes: bytes, filename: str, model: str, instructions: str) -> Tuple[str, LLMUsage]:
    if filename.lower().endswith(".pptx"):
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pdf_bytes))
        text_runs = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_runs.append(f"===Page {slide_num + 1}===\n" + "\n".join(slide_text))
        return "\n\n".join(text_runs), LLMUsage()
    elif filename.lower().endswith(".pdf"):
        import io
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                pages_text = []
                for i, page in enumerate(pdf.pages):
                    page_content = page.extract_text() or ""
                    pages_text.append(f"===Page {i + 1}===\n{page_content}")
                text = "\n\n".join(pages_text)
            return text, LLMUsage()
        except ImportError:
            # Fallback if pdfplumber is not available
            client = PDFClient(pdf_engine="pdf-text", model=model)
            content, _, step_usage = client.send(prompt=instructions, pdf_bytes=pdf_bytes, filename=filename, temperature=0.0, step_name="native_pdf_text")
            return content, step_usage
    else:
        client = PDFClient(pdf_engine="pdf-text", model=model)
        content, _, step_usage = client.send(prompt=instructions, pdf_bytes=pdf_bytes, filename=filename, temperature=0.0, step_name="native_pdf_text")
        return content, step_usage

@task(name="reconcile_texts")
def reconcile_texts(text_content: str, vlm_content: str, text_usage: LLMUsage, vlm_usage: LLMUsage, model: str) -> Tuple[str, LLMUsage]:
    merge_prompt = f"""
You are a specialized VC Technical Analyst merging two transcriptions of the same pitch deck.

TEXT EXTRACTION (Highly accurate numbers, poor layout):
{text_content}

VLM EXTRACTION (Great layout, but may hallucinate numbers):
{vlm_content}

INSTRUCTIONS:
1. Produce a final, unified Markdown transcription.
2. Use the VLM version for the overall structure, tables, and graphs.
3. STRICTLY cross-reference every number, financial metric, and scientific unit against the TEXT EXTRACTION. 
4. If the VLM hallucinated a number, replace it with the correct number from the TEXT EXTRACTION.
5. If there is a discrepancy between the TEXT EXTRACTION and VLM EXTRACTION that you cannot confidently resolve, add a clearly marked section at the bottom of your markdown output exactly named `### ⚠️ Manual Resolution Required`.
   - Format these discrepancies as a concise Markdown table with columns: `Metric`, `VLM Value`, `Native Text Value`.
   - Do NOT flag minor formatting or capitalization differences (e.g. "10 years" vs "10 Years", or "$5m" vs "$5M"). Only flag actual differences in numbers, units, or core meaning.
   - CRITICAL: If there are NO unresolvable discrepancies, DO NOT include this section at all. Do not output the heading and do not state that there are no discrepancies.
6. Output ONLY the final markdown. Do not include your reasoning or introductory text.
"""
    
    client = LLMClient(model=model)
    final_content, final_usage = client.simple_text(prompt=merge_prompt, temperature=0.0, step_name="hybrid_reconciliation")
    
    # Merge all usages
    total_usage = LLMUsage()
    for u in [text_usage, vlm_usage, final_usage]:
        total_usage.total_prompt_tokens += u.total_prompt_tokens
        total_usage.total_completion_tokens += u.total_completion_tokens
        total_usage.total_tokens += u.total_tokens
        total_usage.estimated_cost_usd += u.estimated_cost_usd
        total_usage.call_count += u.call_count
        
    return final_content, total_usage

import json

@flow(name="transcribe_document")
def transcribe_document(pdf_bytes: bytes, filename: str, method: str, model: str) -> Tuple[str, LLMUsage, dict]:
    """
    Process a PDF document using one of four methods:
    'ocr', 'text', 'vlm', or 'vlm_plus_text'.
    Returns (markdown_text, LLMUsage, raw_texts_dict)
    """
    logger.info(json.dumps({
        "event": "transcribe_document_start",
        "filename": filename,
        "method": method,
        "model": model
    }))
    instructions = _load_instructions()
    raw_texts = {}
    
    if method == "ocr":
        res, usage = extract_text_with_ocr(pdf_bytes, filename, model, instructions)
        raw_texts["raw_ocr_text"] = res
    elif method == "paddle_ocr":
        res, usage = extract_text_with_paddle_ocr(pdf_bytes, filename)
        raw_texts["raw_paddle_text"] = res
    elif method == "text":
        res, usage = extract_text_with_native(pdf_bytes, filename, model, instructions)
        raw_texts["raw_native_text"] = res
    elif method == "vlm":
        res, usage = extract_text_with_vlm(pdf_bytes, model, instructions)
        raw_texts["raw_vlm_text"] = res
    elif method == "vlm_plus_text":
        logger.info(json.dumps({
            "event": "transcribe_document_hybrid_start",
            "message": "Running VLM_PLUS_TEXT hybrid extraction... (This will run native text extraction, then VLM extraction, and cross-reference them)"
        }))
        # Run text extraction
        text_content, text_usage = extract_text_with_native(pdf_bytes, filename, model, instructions)
        raw_texts["raw_native_text"] = text_content
        # Run VLM extraction
        vlm_content, vlm_usage = extract_text_with_vlm(pdf_bytes, model, instructions)
        raw_texts["raw_vlm_text"] = vlm_content
        # Cross compare
        res, usage = reconcile_texts(text_content, vlm_content, text_usage, vlm_usage, model)
    elif method == "web_scrape":
        res = pdf_bytes.decode("utf-8")
        usage = LLMUsage()
        raw_texts["raw_native_text"] = res
    else:
        raise ValueError(f"Unknown method: {method}")
        
    logger.info(json.dumps({
        "event": "transcribe_document_complete",
        "status": "success",
        "tokens": usage.total_tokens,
        "cost_usd": round(usage.estimated_cost_usd, 5)
    }))
    return res, usage, raw_texts
